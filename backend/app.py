import os
import sys
import io
from flask import Flask, redirect, url_for, session, request, render_template, jsonify
from google_auth_oauthlib.flow import Flow
from googleapiclient.discovery import build
from langchain_community.vectorstores import FAISS
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain.chains import RetrievalQA
from langchain.text_splitter import CharacterTextSplitter
from langchain_core.documents import Document
import google.generativeai as genai
import logging
from flask_session import Session
from PyPDF2 import PdfReader
from pptx import Presentation
from pathlib import Path
import json

# Load environment variables from .env file
# Load environment variables
try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass  # python-dotenv not installed, using system environment variables

# Allow OAuth over HTTP for development
os.environ['OAUTHLIB_INSECURE_TRANSPORT'] = '1'

# Logging setup

logs_dir = Path(__file__).resolve().parent / "logs"
logs_dir.mkdir(parents=True, exist_ok=True)


logging.basicConfig(
    filename=str(logs_dir / "app.log"),
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s"
)

app = Flask(__name__, 
            template_folder=os.path.abspath('frontend/templates'), 
            static_folder=os.path.abspath('frontend/static'))
app.secret_key = 'rag-chatbot-secret-key-change-in-production-2025'

# Configure server-side sessions
app.config['SESSION_TYPE'] = 'filesystem'
app.config['SESSION_FILE_DIR'] = 'backend/sessions'
app.config['SESSION_PERMANENT'] = True
app.config['PERMANENT_SESSION_LIFETIME'] = 3600  # 1 hour

# Initialize Flask-Session
Session(app)

# Google OAuth setup
CLIENT_SECRETS_FILE = json.loads(os.environ["CLIENT_SECRET_JSON"])
SCOPES = ['https://www.googleapis.com/auth/drive.readonly', 'https://www.googleapis.com/auth/documents.readonly']
API_SERVICE_NAME = 'drive'
API_VERSION = 'v3'

# Environment Variables and AI Service Initialization
GEMINI_API_KEY = os.getenv('GEMINI_API_KEY')
GEMINI_MODEL = os.getenv('GEMINI_MODEL', 'gemini-2.5-flash')

# Initialize Gemini AI services
embeddings = None
llm = None

if GEMINI_API_KEY:
    try:
        # Configure Gemini
        genai.configure(api_key=GEMINI_API_KEY)
        
        # Initialize embeddings and LLM
        embeddings = GoogleGenerativeAIEmbeddings(
            model="models/embedding-001",
            google_api_key=GEMINI_API_KEY
        )
        llm = ChatGoogleGenerativeAI(
            model=GEMINI_MODEL,
            google_api_key=GEMINI_API_KEY,
            temperature=0.7
        )
        logging.info(f"Gemini services initialized successfully with {GEMINI_MODEL}")
    except Exception as e:
        logging.error(f"Error initializing Gemini services: {e}")
        embeddings = None
        llm = None
else:
    logging.warning("No Gemini API Key found. RAG functionality will be limited.")

# Global variables for simplicity (use database in production)
vectorstore = None
selected_docs = []

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/login')
def login():
    try:
        # Clear any existing session
        session.clear()
        
        # Create OAuth flow with explicit scopes
        flow = Flow.from_client_secrets_file(
            CLIENT_SECRETS_FILE, 
            scopes=SCOPES,
            redirect_uri="https://rag-chatbot-zd5g.onrender.com/oauth2callback"
        )
        
        # Request offline access to get refresh token
        authorization_url, state = flow.authorization_url(
            access_type='offline',
            include_granted_scopes='true',
            prompt='consent'  # Force consent screen to ensure refresh token
        )
        
        # Save state in session
        session['state'] = state
        session.permanent = True
        session.modified = True
        
        logging.info(f"Login: Setting state in session: {state}")
        logging.info(f"Login: Session keys after setting: {list(session.keys())}")
        logging.info(f"Login: Session ID: {session.sid if hasattr(session, 'sid') else 'No SID'}")
        logging.info(f"Authorization URL: {authorization_url}")
        logging.info("OAuth flow initiated successfully")
        
        return redirect(authorization_url)
    except FileNotFoundError as e:
        logging.error(f"client_secret.json not found: {e}")
        return "Error: client_secret.json not found. Please download it from Google Cloud Console and place it in the backend/ directory."
    except Exception as e:
        logging.error(f"Error in login: {e}")
        import traceback
        logging.error(f"Traceback: {traceback.format_exc()}")
        return "An error occurred during login setup."

@app.route('/oauth2callback')
def oauth2callback():
    try:
        logging.info(f"Session keys: {list(session.keys())}")
        logging.info(f"Request args: {dict(request.args)}")
        
        # Check for error in response
        if 'error' in request.args:
            error = request.args.get('error')
            logging.error(f"Error returned in OAuth callback: {error}")
            return f"Authentication failed: {error}. Please try again."
        
        if 'code' not in request.args:
            logging.error("No authorization code in response")
            return "Authentication failed: No authorization code received. Please try again."
            
        # Get state from request
        state_param = request.args.get('state')
        logging.info(f"Using state from request: {state_param}")
        
        # Create flow with same scopes and redirect URI
        flow = Flow.from_client_secrets_file(
            CLIENT_SECRETS_FILE, 
            scopes=SCOPES,
            redirect_uri="http://localhost:5000/oauth2callback"
        )
        
        # Get full URL for token exchange
        authorization_response = request.url
        logging.info(f"Authorization response URL: {authorization_response}")
        
        # Exchange code for token
        flow.fetch_token(authorization_response=authorization_response)
        credentials = flow.credentials
        
        # Log credential details (excluding sensitive data)
        logging.info(f"Token received. Expires in: {credentials.expiry}")
        logging.info(f"Scopes granted: {credentials.scopes}")
        logging.info(f"Refresh token received: {'Yes' if credentials.refresh_token else 'No'}")
        
        # Store credentials in session
        session['credentials'] = credentials_to_dict(credentials)
        
        # Make sure the session is saved
        session.modified = True
        
        logging.info("OAuth callback successful")
        logging.info(f"Session after OAuth: {list(session.keys())}")
        return redirect(url_for('docs'))
    except Exception as e:
        logging.error(f"Error in oauth2callback: {e}")
        logging.error(f"Exception type: {type(e)}")
        import traceback
        logging.error(f"Traceback: {traceback.format_exc()}")
        return f"An error occurred during OAuth callback: {str(e)}. Check the logs for details."

@app.route('/docs')
def docs():
    try:
        logging.info(f"Docs route - Session keys: {list(session.keys())}")
        if 'credentials' not in session:
            logging.warning("No credentials in session, redirecting to login")
            return redirect(url_for('login'))
        
        credentials_dict = session['credentials']
        logging.info(f"Credentials dict: {credentials_dict.keys() if credentials_dict else 'None'}")
        
        credentials = dict_to_credentials(credentials_dict)
        
        # Log the full credentials object (except sensitive parts)
        logging.info(f"Token valid: {credentials.valid}")
        logging.info(f"Token expired: {credentials.expired}")
        logging.info(f"Scopes: {credentials.scopes}")
        
        # Make sure we're using the right scopes
        if 'https://www.googleapis.com/auth/drive.readonly' not in credentials.scopes:
            logging.error("Missing required drive.readonly scope")
            return "Error: Missing required drive.readonly scope. Please log in again."
            
        # Build the service with explicit auth
        try:
            logging.info("Building Drive service...")
            logging.info("Using OAuth credentials only for Drive service")
            drive_service = build(
                API_SERVICE_NAME, 
                API_VERSION, 
                credentials=credentials,
                cache_discovery=False
            )
            logging.info("Drive service built successfully")
        except Exception as service_error:
            logging.error(f"Error building Drive service: {service_error}")
            import traceback
            logging.error(f"Service build error traceback: {traceback.format_exc()}")
            return f"Error connecting to Google Drive API: {str(service_error)}"
        
        # Fetch documents from Google Drive with broader query (including Docs, Sheets, PDFs, etc.)
        logging.info("Fetching documents from Drive API with expanded query")
        
        # Query for multiple document types
        query = ("mimeType='application/vnd.google-apps.document' or "
                 "mimeType='application/vnd.google-apps.spreadsheet' or "
                 "mimeType='application/vnd.google-apps.presentation' or "
                 "mimeType='application/pdf' or "
                 "mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation' or "
                 "mimeType='application/vnd.ms-powerpoint' or "
                 "mimeType='text/plain'")
        
        try:
            logging.info(f"Sending Drive API request with query: {query}")
            
            # Create the request
            request = drive_service.files().list(
                q=query,
                pageSize=100,
                fields="nextPageToken, files(id, name, mimeType)"
            )
            
            # Execute with error handling
            results = request.execute()
            
            docs = results.get('files', [])
            logging.info(f"Fetched {len(docs)} documents")
            
        except Exception as api_error:
            logging.error(f"Error querying Drive API: {api_error}")
            import traceback
            logging.error(f"API error traceback: {traceback.format_exc()}")
            
            # Check if this is the specific unregistered caller error
            error_str = str(api_error)
            if "unregistered callers" in error_str or "PERMISSION_DENIED" in error_str:
                logging.error("API key error detected. Check Google Cloud Console configuration.")
                return render_template('api_error.html', 
                    error_message="Google Drive API access error. This application may not be properly configured in Google Cloud Console.",
                    error_details=error_str
                )
            
            return f"Error querying Google Drive: {str(api_error)}"
        
        # Add more detailed logging
        if not docs:
            logging.warning("No documents found in the user's Drive")
            
        # Log the document types found
        mime_types = set(doc.get('mimeType', 'unknown') for doc in docs)
        logging.info(f"Document types found: {mime_types}")
        
        return render_template('docs.html', docs=docs, no_docs=len(docs)==0)
    except Exception as e:
        logging.error(f"Error in docs: {e}")
        import traceback
        logging.error(f"Traceback: {traceback.format_exc()}")
        return f"An error occurred while fetching documents: {str(e)}. Please check the logs for details."

@app.route('/select_docs', methods=['POST'])
def select_docs():
    try:
        global selected_docs, vectorstore
        selected_ids = request.form.getlist('doc_ids')
        selected_docs = selected_ids
        
        if selected_ids and embeddings and llm:
            # Fetch and process documents
            credentials_dict = session['credentials']
            credentials = dict_to_credentials(credentials_dict)
            
            # Build services using OAuth credentials only
            drive_service = build('drive', 'v3', credentials=credentials)
            docs_service = build('docs', 'v1', credentials=credentials)
            
            # First, get document metadata to determine how to process each document
            documents = []
            for doc_id in selected_ids:
                try:
                    # Get document metadata first
                    file_metadata = drive_service.files().get(fileId=doc_id, fields='mimeType,name').execute()
                    mime_type = file_metadata.get('mimeType', '')
                    file_name = file_metadata.get('name', '')
                    
                    logging.info(f"Processing document {file_name} ({doc_id}) with MIME type: {mime_type}")
                    
                    content = ""
                    if mime_type == 'application/vnd.google-apps.document':
                        # Use Google Docs API for Google Docs
                        doc = docs_service.documents().get(documentId=doc_id).execute()
                        content = extract_text_from_doc(doc)
                    elif mime_type == 'text/plain':
                        # Use Drive API to export plain text files
                        response = drive_service.files().get_media(fileId=doc_id).execute()
                        content = response.decode('utf-8')
                    elif mime_type == 'application/pdf':
                        # Use Drive API to get PDF content and extract text
                        response = drive_service.files().get_media(fileId=doc_id).execute()
                        content = extract_text_from_pdf(response)
                        if not content:
                            logging.warning(f"Could not extract text from PDF {file_name}")
                            continue
                    elif mime_type == 'application/vnd.google-apps.presentation':
                        # Export Google Slides as PPTX and extract text
                        response = drive_service.files().export_media(
                            fileId=doc_id, 
                            mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
                        ).execute()
                        content = extract_text_from_pptx(response)
                        if not content:
                            logging.warning(f"Could not extract text from Google Slides {file_name}")
                            continue
                    elif mime_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation', 'application/vnd.ms-powerpoint']:
                        # Use Drive API to get PowerPoint content and extract text
                        response = drive_service.files().get_media(fileId=doc_id).execute()
                        content = extract_text_from_pptx(response)
                        if not content:
                            logging.warning(f"Could not extract text from PowerPoint {file_name}")
                            continue
                    else:
                        logging.warning(f"Unsupported document type: {mime_type} for {file_name}")
                        continue
                    
                    if content.strip():
                        documents.append(Document(page_content=content, metadata={'source': doc_id, 'name': file_name}))
                        logging.info(f"Successfully processed document: {file_name}")
                    else:
                        logging.warning(f"Document {file_name} appears to be empty")
                        
                except Exception as doc_error:
                    logging.error(f"Error processing document {doc_id}: {doc_error}")
                    continue
            
            # Split and embed
            if documents:
                text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=0)
                docs = text_splitter.split_documents(documents)
                vectorstore = FAISS.from_documents(docs, embeddings)
                logging.info(f"Processed {len(documents)} documents into vectorstore")
            else:
                logging.warning("No documents were successfully processed")
                vectorstore = None
                return "Error: No documents could be processed. Please check that the selected documents are accessible and properly formatted."
        elif selected_ids:
            vectorstore = None  # API key not set
            logging.warning("Documents selected but OpenAI API key not set")
        
        return redirect(url_for('chat'))
    except Exception as e:
        logging.error(f"Error in select_docs: {e}")
        return f"An error occurred while selecting documents: {str(e)}"

@app.route('/chat')
def chat():
    return render_template('chat.html', 
                         gemini_available=bool(GEMINI_API_KEY and embeddings and llm),
                         selected_docs_count=len(selected_docs))

@app.route('/logout')
def logout():
    # Clear the session
    session.clear()
    logging.info("User logged out, session cleared")
    return redirect(url_for('index'))

@app.route('/query', methods=['POST'])
def query():
    try:
        user_query = request.json['query']
        logging.info(f"Processing query: {user_query[:50]}...")
        
        if not GEMINI_API_KEY:
            response = "⚠ Gemini API key is not configured. Please set the GEMINI_API_KEY environment variable to enable RAG functionality with Gemini."
        elif vectorstore and llm:
            logging.info("Using RAG pipeline with selected documents and Gemini")
            qa_chain = RetrievalQA.from_chain_type(llm=llm, chain_type="stuff", retriever=vectorstore.as_retriever())
            result = qa_chain.invoke({"query": user_query})
            
            # Extract the result text
            if isinstance(result, dict):
                result_text = result.get('result', str(result))
            else:
                result_text = str(result)
            
            # Check if answer is relevant (simple check)
            if "I don't know" in result_text or len(result_text.strip()) < 50:
                general = generate_general_response(user_query)
                response = f"I couldn't find the answer in your selected documents. However, based on general knowledge: {general}"
            else:
                response = result_text
        elif selected_docs and not vectorstore:
            response = f"Documents are selected ({len(selected_docs)} docs) but the RAG pipeline couldn't be initialized. This might be due to missing Google API key or other configuration issues."
        elif selected_docs:
            response = f"Documents selected but vectorstore not ready. Falling back to general knowledge: {generate_general_response(user_query)}"
        else:
            if GEMINI_API_KEY:
                response = f"No documents selected. Using general knowledge: {generate_general_response(user_query)}"
            else:
                response = "No documents selected and Gemini API key not configured. Please select documents and configure your Gemini API key to use the chatbot."
        
        logging.info(f"Query processed successfully")
        return jsonify({'response': response})
    except Exception as e:
        logging.error(f"Error in query: {e}")
        import traceback
        logging.error(f"Query error traceback: {traceback.format_exc()}")
        return jsonify({'response': f"An error occurred: {str(e)}"})

def credentials_to_dict(credentials):
    return {
        'token': credentials.token,
        'refresh_token': credentials.refresh_token,
        'token_uri': credentials.token_uri,
        'client_id': credentials.client_id,
        'client_secret': credentials.client_secret,
        'scopes': credentials.scopes
    }
    
def dict_to_credentials(credentials_dict):
    from google.oauth2.credentials import Credentials
    return Credentials(
        token=credentials_dict['token'],
        refresh_token=credentials_dict['refresh_token'],
        token_uri=credentials_dict['token_uri'],
        client_id=credentials_dict['client_id'],
        client_secret=credentials_dict['client_secret'],
        scopes=credentials_dict['scopes']
    )

def extract_text_from_doc(doc):
    content = ""
    for element in doc.get('body', {}).get('content', []):
        if 'paragraph' in element:
            for para_element in element['paragraph']['elements']:
                if 'textRun' in para_element:
                    content += para_element['textRun']['content']
    return content

def extract_text_from_pdf(pdf_data):
    """Extract text from PDF data using PyPDF2"""
    try:
        pdf_stream = io.BytesIO(pdf_data)
        pdf_reader = PdfReader(pdf_stream)
        content = ""
        
        for page in pdf_reader.pages:
            content += page.extract_text() + "\n"
        
        return content.strip()
    except Exception as e:
        logging.error(f"Error extracting text from PDF: {e}")
        return ""

def extract_text_from_pptx(pptx_data):
    """Extract text from PowerPoint data using python-pptx"""
    try:
        pptx_stream = io.BytesIO(pptx_data)
        presentation = Presentation(pptx_stream)
        content = ""
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            content += f"\n--- Slide {slide_num} ---\n"
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content += shape.text + "\n"
                
                # Extract text from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            content += " | ".join(row_text) + "\n"
            
            content += "\n"
        
        return content.strip()
    except Exception as e:
        logging.error(f"Error extracting text from PowerPoint: {e}")
        return ""

def generate_general_response(query):
    try:
        if not GEMINI_API_KEY:
            return "Gemini API key not set. Cannot generate general response."
        
        # Simple general response using Gemini
        model = genai.GenerativeModel(GEMINI_MODEL)
        response = model.generate_content(query)
        return response.text.strip()
    except Exception as e:
        logging.error(f"Error in generate_general_response: {e}")
        return "Error generating general response."

if __name__ == '__main__':
    app.run(debug=True)
